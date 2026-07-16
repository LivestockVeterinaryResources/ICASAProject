import os
from pptx import Presentation
from pptx.util import Emu

src = r"C:\Users\steve\Documents\R\ICASAProject\Apley ICASA 4-2-2026.pptx"
outdir = r"C:\Users\steve\AppData\Local\Temp\claude\C--Users-steve-Documents-R-ICASAProject\eab90873-9b28-41b6-a7fc-2b2be4934377\scratchpad\apley_media"
os.makedirs(outdir, exist_ok=True)

prs = Presentation(src)
sw, sh = prs.slide_width, prs.slide_height
seen = {}
report = []
for i, slide in enumerate(prs.slides, 1):
    pics = []
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            texts.append(shape.text_frame.text.strip().replace("\n"," | "))
        if shape.shape_type == 13 or getattr(shape, "image", None) is not None:  # PICTURE
            try:
                img = shape.image
            except Exception:
                continue
            # fraction of slide area
            try:
                area_frac = (shape.width * shape.height) / (sw * sh)
            except Exception:
                area_frac = 0
            ext = img.ext
            key = img.sha1
            if key not in seen:
                fn = f"img_{key[:8]}.{ext}"
                with open(os.path.join(outdir, fn), "wb") as f:
                    f.write(img.blob)
                seen[key] = fn
            pics.append((seen[key], round(area_frac,2), img.size))
    report.append((i, texts, pics))

for i, texts, pics in report:
    big = [p for p in pics if p[1] >= 0.15]
    if big or texts:
        t = (texts[0][:60] if texts else "(no title text)")
        print(f"S{i:>2} | {t}")
        for fn, frac, size in pics:
            flag = "  <== BIG FIGURE" if frac >= 0.15 else ""
            print(f"       {fn}  area={frac} px={size}{flag}")
print("\nUnique images:", len(seen))
